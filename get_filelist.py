#pip install baidu-aip
from aip import AipOcr
#pip install xlrd
#import xlrd
from xlrd import open_workbook
import os
from os.path import isfile, isdir, join, exists
from os import listdir, remove
#pip install csv23
#import csv23
from csv23 import open_csv
#pip install docx2txt
from docx2txt import process
#pip install pypiwin32
#from win32com import client as wc
#pip install python-pptx
from pptx import Presentation
#pip install wand
#from wand.image import Image
# 使用 wand 异常，缺少 ImageMagick 支持
# http://docs.wand-py.org/en/latest/guide/install.html#install-imagemagick-on-windows
# https://imagemagick.org/script/download.php#windows
# 使用 wand 异常，FailedToExecuteCommand `"gswin32c.exe"
# http://ghostscript.com/download/gsdnld.html

def get_filelist(dir, Filelist):
    newDir = dir

    # if os.path.isfile(dir):
    if isfile(dir):

        Filelist.append(dir)

    # 若只是要返回文件文，使用这个

    # Filelist.append(os.path.basename(dir))

    # elif os.path.isdir(dir):
    elif isdir(dir):

        # for s in os.listdir(dir):
        for s in listdir(dir):
            # 如果需要忽略某些文件夹，使用以下代码
            # if s == "xxx":
            # continue

            # newDir = os.path.join(dir, s)
            newDir = join(dir, s)

            get_filelist(newDir, Filelist)

    return Filelist


def readTxt(fileUrl):
    content = ""
    # if os.path.exists(fileUrl):
    if exists(fileUrl):
        with open(fileUrl, 'r', encoding='utf-8') as f:
            for l in f:
                temp = l.rstrip('\n').rstrip().split('\t')[0]
                # print(temp)
                content += temp.replace(' ',' ')
    return content

def readDocx(fileUrl):
    content = ""
    # if os.path.exists(fileUrl):
    if exists(fileUrl):
        # content = docx2txt.process(fileUrl)
        content = process(fileUrl)
        #content = "".join(content.split())
    return content

# 读取 doc 文件，安装 pypiwin32，操作本地word程序，将doc 转为docx，再调用读取 docx 文件方法
# def readDoc(fileUrl):
#     AbsolutePath = os.path.abspath(fileUrl)
#     word = wc.Dispatch('Word.Application')
#     doc = word.Documents.Open(AbsolutePath)
#     # 保存临时文件
#     doc.SaveAs(AbsolutePath + ".docx", 12, False, "", True, "", False, False, False, False) # 转化后路径下的文件
#     doc.Close()
#     word.Quit()
#     content = readDocx(fileUrl + ".docx")
#     # 移除临时文件
#     os.remove(fileUrl + ".docx")
#     return content

# 读取 csv 文件，返回文件内容，默认utf-8，如果解析不了，使用gbk解析
def readCsv(fileUrl):
    content = ""
    # if os.path.exists(fileUrl):
    if exists(fileUrl):
        try:
            with open_csv(fileUrl) as reader:
                for row in reader:
                    content += (''.join(row)).replace(' ','')
        except Exception as e:
            with open_csv(fileUrl, encoding='gbk') as reader:
                for row in reader:
                    content += (''.join(row)).replace(' ','')
    return content

# 读取 图片 文件，返回文件内容
def readImage(fileUrl):
    content = ""
    if exists(fileUrl):
        APP_ID = 'xxxxx'
        API_KEY = 'xxxxxxxxxxxxxxxx'
        SECRET_KEY = 'xxxxxxxxxxxxxxxxxxxxxxx'
        client = AipOcr(APP_ID, API_KEY, SECRET_KEY)
        with open(fileUrl,'rb') as f:
            img = f.read()
            msg = client.basicGeneral(img)
            for i in msg.get('words_result'):
                temp = i.get('words')
                content += temp.replace(' ','')
    return content

# 读取 pptx 文件 ,默认读取正文，默认不读取表格，默认不读取图片
def readPptx(fileUrl,extend_table = False,extend_image = False):
    content = ""
    ppt = Presentation(fileUrl)

    for slide in ppt.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                # 提取图片文字
                if extend_image and hasattr(shape,'image'):
                    # 图片存储本地
                    with open(shape.image.filename, 'wb') as f:
                        f.write(shape.image.blob)
                        f.close()
                    # 调用图片文字识别
                    content += readImage(shape.image.filename)
                    # 移除临时图片
                    remove(shape.image.filename)
                # 提取表格内容
                if extend_table and shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            content += cell.text
            else:
                content += shape.text

    content = "".join(content.split())
    return content

# 读取 ppt 文件，安装 pypiwin32，操作本地ppt程序，将ppt 转为pptx，再调用读取 pptx 文件方法
# def readPpt(fileUrl,extend_table = False,extend_image = False):
#     AbsolutePath = os.path.abspath(fileUrl)
#     powerpoint = wc.Dispatch('PowerPoint.Application')
#     ppt = powerpoint.Presentations.Open(AbsolutePath)
#     # 保存临时文件
#     ppt.SaveAs(AbsolutePath + ".pptx")
#     powerpoint.Quit()
#     content = readPptx(fileUrl + ".pptx",extend_table,extend_image)
#     # 移除临时文件
#     os.remove(fileUrl + ".pptx")
#     return content


# 读取 pdf 文件
# def readPdf(fileUrl):
#     content = ""
#     # 将pdf文件转为jpg图片文件
#     # ./PDF_FILE_NAME 为pdf文件路径和名称
#     image_pdf = Image(filename=fileUrl, resolution=300)
#     image_jpeg = image_pdf.convert('jpg')
#
#     # wand已经将PDF中所有的独立页面都转成了独立的二进制图像对象。我们可以遍历这个大对象，并把它们加入到req_image序列中去。
#     req_image = []
#     for img in image_jpeg.sequence:
#         img_page = Image(image=img)
#         req_image.append(img_page.make_blob('jpg'))
#
#     # 遍历req_image,保存为图片文件
#
#     for img in req_image:
#         ff = open(fileUrl + '.jpg', 'wb')
#         ff.write(img)
#         ff.close()
#         # 调用图片文字识别
#         content += readImage(fileUrl + '.jpg')
#         # 移除临时图片
#         os.remove(fileUrl + '.jpg')
#     return content

# 读取 excel 文件，返回文件内容
def readExcel(fileUrl):
    content = ""
    if exists(fileUrl):
        excelfile = open_workbook(fileUrl)
        for name in excelfile.sheet_names():
            sheet = excelfile.sheet_by_name(name)
            sheet_rows = sheet.nrows
            sheet_cols = sheet.ncols
            for rowi in range(sheet_rows):
                temp = sheet.row_values(rowi)
                content += (''.join(map(str,temp))).replace(' ','')
    return content